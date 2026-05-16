"""Document ingestion — reads PDFs, DOCX, PPTX, CSV, JSON, images, plain
text, and GitHub repos into page/file dicts.  Also supports web URLs."""

import csv
import hashlib
import json as json_mod
import logging
import re
import subprocess
import tempfile
from pathlib import Path

from pypdf import PdfReader

log = logging.getLogger(__name__)

# File extensions to include when reading a repo
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".rb",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".swift", ".kt", ".scala",
    ".sh", ".bash", ".zsh", ".yaml", ".yml", ".toml", ".json",
    ".html", ".css", ".scss", ".sql", ".r", ".R", ".md", ".txt",
    ".cfg", ".ini", ".env.example", ".dockerfile", ".tf",
}

# Directories to always skip
SKIP_DIRS = {
    ".git", "node_modules", "__pycache__", ".venv", "venv",
    "dist", "build", ".eggs", ".tox", ".mypy_cache",
}

GITHUB_URL_RE = re.compile(
    r"^https?://github\.com/(?P<owner>[^/]+)/(?P<repo>[^/]+?)(?:\.git)?/?$"
)

URL_RE = re.compile(r"^https?://")
S3_URI_RE = re.compile(r"^s3://(?P<bucket>[^/]+)/(?P<key>.+)$")

PLAIN_TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".log"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif"}


def is_s3_uri(source: str) -> bool:
    return bool(S3_URI_RE.match(source))


def parse_s3_uri(source: str) -> tuple[str, str]:
    m = S3_URI_RE.match(source)
    if not m:
        raise ValueError(f"Not an S3 URI: {source}")
    return m["bucket"], m["key"]


# ── Individual readers ────────────────────────────────────────────────


def read_pdf(file_path: str) -> list[dict]:
    """Extract text from each page of a PDF, returning page dicts."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PDF not found: {file_path}")

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append({"text": text, "page": i + 1})
        else:
            log.warning("Page %d is empty or unreadable", i + 1)
    return pages


def read_plain_text(file_path: str) -> list[dict]:
    """Read a plain-text or Markdown file as a single page."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []
    return [{"text": text, "page": 1}]


def read_docx(file_path: str) -> list[dict]:
    """Read a DOCX file, one page-dict per paragraph group."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required for DOCX support: pip install python-docx")
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"DOCX not found: {file_path}")
    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs).strip()
    if not text:
        return []
    return [{"text": text, "page": 1}]


def read_pptx(file_path: str) -> list[dict]:
    """Read a PPTX file, one page-dict per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPTX support: pip install python-pptx")
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {file_path}")
    prs = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        text = "\n".join(parts).strip()
        if text:
            pages.append({"text": text, "page": i + 1})
    return pages


def read_csv(file_path: str) -> list[dict]:
    """Read a CSV file, converting each row to text."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"CSV not found: {file_path}")
    text_parts = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.DictReader(f)
        for row in reader:
            line = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
            if line:
                text_parts.append(line)
    if not text_parts:
        return []
    # Return as a single page — the chunker will split it
    return [{"text": "\n".join(text_parts), "page": 1}]


def read_json(file_path: str) -> list[dict]:
    """Read a JSON file, converting to readable text."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"JSON not found: {file_path}")
    raw = path.read_text(encoding="utf-8", errors="ignore")
    data = json_mod.loads(raw)
    # Pretty-print so the chunker can split on newlines
    text = json_mod.dumps(data, indent=2, ensure_ascii=False, default=str)
    if not text.strip():
        return []
    return [{"text": text, "page": 1}]


def read_image(file_path: str) -> list[dict]:
    """OCR an image file using pytesseract."""
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        raise ImportError(
            "Pillow and pytesseract are required for image OCR: "
            "pip install Pillow pytesseract  (and install Tesseract binary)"
        )
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {file_path}")
    img = Image.open(path)
    text = pytesseract.image_to_string(img).strip()
    if not text:
        log.warning("OCR returned no text for %s", file_path)
        return []
    return [{"text": text, "page": 1}]


def read_url(url: str) -> list[dict]:
    """Fetch a web page and extract its main text content."""
    try:
        import requests
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError(
            "requests and beautifulsoup4 are required for URL support: "
            "pip install requests beautifulsoup4"
        )
    resp = requests.get(url, timeout=30, headers={"User-Agent": "ProfessionalRAG/1.0"})
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    # Remove script/style elements
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    if not text:
        return []
    return [{"text": text, "page": url}]


def is_github_url(source: str) -> bool:
    """Check if a source string is a GitHub repo URL."""
    return bool(GITHUB_URL_RE.match(source))


# ── Dispatcher ────────────────────────────────────────────────────────


def read_source(source: str) -> list[dict]:
    """Auto-detect the source type and read it into page dicts.

    Accepts: PDF, DOCX, PPTX, CSV, JSON, plain text, Markdown, images,
    GitHub URLs, local directories, or web URLs.

    Returns a list of dicts with keys: text, page, source.
    """
    # S3 object
    if is_s3_uri(source):
        return read_s3(source)

    # GitHub repo
    if is_github_url(source):
        return read_repo(source)

    # Web URL (non-GitHub)
    if URL_RE.match(source):
        pages = read_url(source)
        for p in pages:
            p.setdefault("source", source)
        return pages

    # Local directory → repo reader
    path = Path(source)
    if path.is_dir():
        return read_repo(source)

    # Local file — dispatch by extension
    if not path.exists():
        raise FileNotFoundError(f"Source not found: {source}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        pages = read_pdf(source)
    elif suffix == ".docx":
        pages = read_docx(source)
    elif suffix == ".pptx":
        pages = read_pptx(source)
    elif suffix == ".csv":
        pages = read_csv(source)
    elif suffix == ".json":
        pages = read_json(source)
    elif suffix in IMAGE_EXTENSIONS:
        pages = read_image(source)
    elif suffix in PLAIN_TEXT_EXTENSIONS:
        pages = read_plain_text(source)
    else:
        # Fallback: try reading as plain text
        log.info("Unknown extension %s — trying as plain text", suffix)
        pages = read_plain_text(source)

    for p in pages:
        p.setdefault("source", source)
    return pages


def read_repo(source: str) -> list[dict]:
    """Read source files from a GitHub URL or local directory path.

    Returns a list of dicts with keys: text, page (file path), source.
    """
    if is_github_url(source):
        repo_dir = _clone_repo(source)
        cleanup = True
    else:
        repo_dir = Path(source)
        if not repo_dir.is_dir():
            raise FileNotFoundError(f"Directory not found: {source}")
        cleanup = False

    try:
        return _read_directory(repo_dir, source)
    finally:
        if cleanup:
            import shutil
            shutil.rmtree(repo_dir, ignore_errors=True)


def _clone_repo(url: str) -> Path:
    """Shallow-clone a GitHub repo into a temp directory."""
    tmp = Path(tempfile.mkdtemp(prefix="rag_repo_"))
    log.info("Cloning %s → %s", url, tmp)
    subprocess.run(
        ["git", "clone", "--depth", "1", url, str(tmp)],
        check=True,
        capture_output=True,
        text=True,
    )
    return tmp


def _read_directory(root: Path, source: str) -> list[dict]:
    """Walk a directory and read code files into page dicts."""
    docs = []
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        # Skip hidden/unwanted dirs
        if any(skip in path.parts for skip in SKIP_DIRS):
            continue
        if path.suffix not in CODE_EXTENSIONS:
            continue
        try:
            text = path.read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            log.warning("Could not read %s", path)
            continue
        if not text:
            continue

        rel_path = str(path.relative_to(root))
        docs.append({
            "text": f"# File: {rel_path}\n\n{text}",
            "page": rel_path,
            "source": source,
        })
    log.info("Read %d files from %s", len(docs), source)
    return docs


def read_s3(source: str) -> list[dict]:
    """Download an S3 object to a temp file and dispatch by extension.

    Uses the local boto3 default credential chain — on EC2 this is the
    instance role; locally it's whatever ~/.aws/credentials or env vars
    supply. The temp file is cleaned up after parsing.
    """
    import boto3  # imported lazily so tests can stub without the module loaded

    bucket, key = parse_s3_uri(source)
    suffix = Path(key).suffix.lower() or ".bin"

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp_path = tmp.name

    try:
        log.info("Downloading s3://%s/%s → %s", bucket, key, tmp_path)
        boto3.client("s3").download_file(bucket, key, tmp_path)
        pages = read_source(tmp_path)  # recursive dispatch on the local copy
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    # Replace the temp-file source path with the original S3 URI so chunks
    # remember where they came from, not the ephemeral local path.
    for p in pages:
        p["source"] = source
    return pages


def s3_fingerprint(source: str) -> str:
    """ETag of the S3 object — changes whenever the object's bytes do.

    For single-part uploads ETag is the MD5; for multi-part it's a
    deterministic composite. Either way it's stable per content version,
    which is all the fingerprint cache needs.
    """
    import boto3

    bucket, key = parse_s3_uri(source)
    head = boto3.client("s3").head_object(Bucket=bucket, Key=key)
    # ETag is wrapped in quotes; strip them for a clean hex string
    etag = head["ETag"].strip('"')
    return f"s3://{bucket}/{key}@{etag}"


def file_fingerprint(file_path: str) -> str:
    """SHA-256 hash of file contents — used to detect changes for re-ingestion."""
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()


def dir_fingerprint(source: str) -> str:
    """SHA-256 hash based on the source identifier (URL or path)."""
    return hashlib.sha256(source.encode()).hexdigest()
